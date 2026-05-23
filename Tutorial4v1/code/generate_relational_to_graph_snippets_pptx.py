from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "outputs" / "relational_to_graph_snippets.pptx"

SLIDE_1_CODE = """import pandas as pd
import networkx as nx

customers = pd.read_csv("data/customers.csv")
transactions = pd.read_csv("data/transactions.csv")
interactions = pd.read_csv("data/interactions.csv")

G = nx.Graph()

# Customer rows -> customer nodes
for row in customers.itertuples(index=False):
    G.add_node(
        row.customer_id,
        node_type="customer",
        region=row.region,
        channel_pref=row.channel_pref,
    )
"""

SLIDE_2_CODE = """# Transaction rows -> transaction nodes (+ edge to owner customer)
for i, row in enumerate(transactions.itertuples(index=False), start=1):
    t = f"txn_{i:03d}"
    G.add_node(t, node_type="transaction", category=row.category, amount=row.amount)
    G.add_edge(t, row.customer_id, relation="made_transaction")

# Interaction rows -> interaction nodes (+ 2 endpoint edges)
for i, row in enumerate(interactions.itertuples(index=False), start=1):
    e = f"int_{i:03d}"
    G.add_node(e, node_type="interaction")
    G.add_edge(e, row.customer_i, relation="interaction_endpoint")
    G.add_edge(e, row.customer_j, relation="interaction_endpoint")

print(G.number_of_nodes(), G.number_of_edges())  # expected: 106, 116
"""


def add_code_slide(prs: Presentation, title: str, subtitle: str, code: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.5))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_run = title_tf.paragraphs[0].runs[0]
    title_run.font.size = Pt(30)
    title_run.font.bold = True
    title_run.font.name = "Calibri"

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_run = subtitle_tf.paragraphs[0].runs[0]
    subtitle_run.font.size = Pt(16)
    subtitle_run.font.name = "Calibri"
    subtitle_run.font.color.rgb = RGBColor(80, 80, 80)

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.55))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(245, 247, 250)
    code_box.line.color.rgb = RGBColor(215, 220, 230)

    code_tf = code_box.text_frame
    code_tf.clear()
    code_tf.word_wrap = True
    code_tf.margin_left = Inches(0.18)
    code_tf.margin_right = Inches(0.18)
    code_tf.margin_top = Inches(0.12)
    code_tf.margin_bottom = Inches(0.12)

    para = code_tf.paragraphs[0]
    para.text = code
    para.alignment = PP_ALIGN.LEFT
    run = para.runs[0]
    run.font.name = "Consolas"
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(20, 20, 20)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_code_slide(
        prs,
        title="Relational Tables -> Graph: Setup + Customer Nodes",
        subtitle="Recitation 4 dataset (customers, transactions, interactions)",
        code=SLIDE_1_CODE,
    )
    add_code_slide(
        prs,
        title="Relational Tables -> Graph: Transaction + Interaction Rows",
        subtitle="Each row becomes a node; edges encode ownership/endpoints",
        code=SLIDE_2_CODE,
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
